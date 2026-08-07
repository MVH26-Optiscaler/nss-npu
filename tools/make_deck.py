#!/usr/bin/env python3
"""Builds the hackathon deck: neural upscaling on Snapdragon, three slides."""
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x0E, 0x11, 0x17)
FG = RGBColor(0xF2, 0xF4, 0xF8)
MUTED = RGBColor(0x9A, 0xA4, 0xB2)
ACCENT = RGBColor(0x7C, 0x4D, 0xFF)
GOOD = RGBColor(0x00, 0xD9, 0xA3)
W, H = Inches(13.333), Inches(7.5)


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, size=18, color=FG, bold=False, space=10, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(runs if isinstance(runs, list) else [runs]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if isinstance(item, tuple):
            body, opts = item
        else:
            body, opts = item, {}
        r = p.add_run()
        r.text = body
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
        r.font.name = "Inter"
    return box


def rule(s, y, x=Inches(0.9), w=Inches(11.5), color=ACCENT, h=Emu(28575)):
    bar = s.shapes.add_shape(1, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def metric(s, x, y, value, label, color=GOOD):
    text(s, x, y, Inches(3.0), Inches(0.7), value, size=34, bold=True, color=color, space=2)
    text(s, x, y + Inches(0.62), Inches(3.0), Inches(0.6), label, size=12, color=MUTED)


def build(path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------------------------------------------------------------- slide 1
    s = slide(prs)
    text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5),
         "WINDOWS GAMES ON SNAPDRAGON", size=13, color=ACCENT, bold=True)
    text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         "Neural upscaling, running on the NPU", size=40, bold=True)
    rule(s, Inches(1.95))
    text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.6),
         "Galaxy S25 Ultra  ·  Snapdragon 8 Elite  ·  Adreno 830  ·  Hexagon HTP V79",
         size=15, color=MUTED)

    text(s, Inches(0.9), Inches(3.0), Inches(5.4), Inches(3.4), [
        ("What we did", {"size": 20, "bold": True, "color": ACCENT}),
        "Replaced Intel's libxess.dll to intercept the game's XeSS calls — giving us "
        "its full G-buffer: colour, depth, motion vectors and jitter.",
        "That is the one interception point that survives ARM64EC translation, where "
        "runtime code patching does not.",
        "From there the upscaler is ours to choose.",
    ], size=15, space=12)

    text(s, Inches(6.9), Inches(3.0), Inches(5.5), Inches(3.4), [
        ("Shipping today", {"size": 20, "bold": True, "color": ACCENT}),
        "SGSR 1.0 — Half-Life 2, D3D9 through DXVK. 640×480 upscaled to 1280×720, "
        "full screen.",
        "SGSR 2.0 — The Witcher 3, D3D12 through vkd3d. Ported GLSL to HLSL; both "
        "passes cost ~1 ms of frame time.",
        "Arm NSS — running on the Hexagon NPU, inside the game's own process.",
    ], size=15, space=12)

    metric(s, Inches(0.9), Inches(6.35), "14–25 fps", "The Witcher 3, upscaled")
    metric(s, Inches(4.2), Inches(6.35), "7.5 ms", "NSS inference on the NPU")
    metric(s, Inches(7.5), Inches(6.35), "~1 ms", "SGSR 2 GPU cost")
    metric(s, Inches(10.5), Inches(6.35), "0 code", "changes to the game")

    # ---------------------------------------------------------------- slide 2
    s = slide(prs)
    text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5),
         "GETTING NSS ONTO THE NPU", size=13, color=ACCENT, bold=True)
    text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         "Four walls, and what was behind each", size=40, bold=True)
    rule(s, Inches(1.95))

    text(s, Inches(0.9), Inches(2.25), Inches(5.4), Inches(4.2), [
        ("The model does not ship as ONNX", {"size": 17, "bold": True, "color": GOOD}),
        "Arm publishes NSS as PyTorch weights and a compiled VGF. We rebuilt it from "
        "their model gym and validated the export against Arm's own reference tensors "
        "before it ever reached the device — 0.9994 correlation.",
        ("Quantisation is not a checkbox", {"size": 17, "bold": True, "color": GOOD}),
        "int8 per-channel is mandatory: per-tensor quietly collapsed one output head "
        "to 0.26 correlation while the other still read 0.996. A single healthy metric "
        "proves nothing.",
    ], size=14, space=10)

    text(s, Inches(6.9), Inches(2.25), Inches(5.5), Inches(4.2), [
        ("The NPU is not simply open", {"size": 17, "bold": True, "color": GOOD}),
        "An app cannot open the DSP node directly, unsigned process domains must be "
        "enabled explicitly, and wine runs in a linker namespace where the vendor "
        "FastRPC libraries are unreachable until staged under /data.",
        ("Failure is silent", {"size": 17, "bold": True, "color": GOOD}),
        "When QNN cannot reach the DSP, ONNX Runtime places every node on the CPU and "
        "reports success. Three measurements read as NPU results while running on the "
        "CPU. The bridge now refuses to start rather than fall back.",
    ], size=14, space=10)

    rule(s, Inches(6.45), color=RGBColor(0x2A, 0x2F, 0x3A))
    text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.6),
         [("Layout was worth 2.6×.  ", {"bold": True, "color": GOOD}),
          "The HTP is NHWC-native — an NCHW interface made QNN convert layout around "
          "every operator. 19.8 ms → 7.5 ms, with the graph unchanged."],
         size=14, color=MUTED)

    # ---------------------------------------------------------------- slide 3
    s = slide(prs)
    text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5),
         "WHY THIS GENERALISES", size=13, color=ACCENT, bold=True)
    text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         "One hook, any game, every Snapdragon", size=40, bold=True)
    rule(s, Inches(1.95))

    text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2),
         "We did not modify The Witcher 3 or Half-Life 2. We replaced the upscaler "
         "library the game already loads. Every title shipping XeSS, FSR or DLSS "
         "exposes the same surface — colour, depth, motion vectors, jitter — which is "
         "exactly what a neural upscaler needs.",
         size=17, space=6)

    text(s, Inches(0.9), Inches(3.7), Inches(3.5), Inches(2.4), [
        ("Any game", {"size": 18, "bold": True, "color": ACCENT}),
        "The interception point is an industry API, not a per-title mod. The work "
        "carries to the next game for free.",
    ], size=14, space=8)
    text(s, Inches(4.9), Inches(3.7), Inches(3.5), Inches(2.4), [
        ("Any Snapdragon", {"size": 18, "bold": True, "color": ACCENT}),
        "Adreno for the shaders, Hexagon for the network. Silicon that ships in every "
        "flagship and sits idle while you game.",
    ], size=14, space=8)
    text(s, Inches(8.9), Inches(3.7), Inches(3.5), Inches(2.4), [
        ("Any frame budget", {"size": 18, "bold": True, "color": ACCENT}),
        "Render at half resolution, reconstruct on the NPU. The GPU keeps the "
        "milliseconds it would have spent upscaling.",
    ], size=14, space=8)

    rule(s, Inches(6.2), color=RGBColor(0x2A, 0x2F, 0x3A))
    text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.9),
         [("Where this goes:  ", {"bold": True, "color": GOOD}),
          "a Linux distribution for Snapdragon that runs the Windows game library "
          "natively-fast — Proton and vkd3d tuned for Adreno, with NPU upscaling and "
          "frame reconstruction built into the compositor rather than bolted onto each "
          "title."],
         size=15, color=MUTED)

    prs.save(path)
    print(f"wrote {path}")


if __name__ == "__main__":
    out = pathlib.Path.home() / "Desktop" / "snapdragon-neural-upscaling.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    build(str(out))
